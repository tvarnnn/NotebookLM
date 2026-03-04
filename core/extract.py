from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional, Tuple

import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation


@dataclass(frozen=True)
class ExtractedText:
    text: str
    # used for citations later
    location_hint: str


def extract_txt(path: Path) -> ExtractedText:
    txt = path.read_text(encoding="utf-8", errors="ignore")
    return ExtractedText(text=txt, location_hint="Text")


def extract_pdf(path: Path) -> ExtractedText:
    reader = PdfReader(str(path))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        # mark page boundaries for citations
        parts.append(f"\n[PAGE {i}]\n{page_text}\n")
    return ExtractedText(text="\n".join(parts).strip(), location_hint="Pages")


def extract_pptx(path: Path) -> ExtractedText:
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text_parts.append(shape.text)
        slide_text = "\n".join(slide_text_parts).strip()
        parts.append(f"\n[SLIDE {i}]\n{slide_text}\n")
    return ExtractedText(text="\n".join(parts).strip(), location_hint="Slides")


def extract_url(url: str, timeout_s: int = 20) -> ExtractedText:
    r = requests.get(url, timeout=timeout_s, headers={"User-Agent": "Mozilla/5.0"})
    r.raise_for_status()
    soup = BeautifulSoup(r.text, "html.parser")

    # remove junk
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()

    text = soup.get_text(separator="\n")
    # basic cleanup
    lines = [ln.strip() for ln in text.splitlines()]
    lines = [ln for ln in lines if ln]
    cleaned = "\n".join(lines)
    return ExtractedText(text=cleaned, location_hint="URL")


def extract_any(path: Path) -> Optional[ExtractedText]:
    ext = path.suffix.lower()
    if ext == ".txt":
        return extract_txt(path)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    return None