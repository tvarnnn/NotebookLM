from __future__ import annotations

import json
import shutil
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import List

from pypdf import PdfReader
from pptx import Presentation

from storage.paths import notebook_root


@dataclass
class SourceMeta:
    id: str           # filename (unique per notebook)
    type: str         # pdf/txt/pptx/url
    enabled: bool
    raw_path: str
    text_path: str


def _sources_json_path(nb_root: Path) -> Path:
    return nb_root / "sources.json"


def _load_sources(nb_root: Path) -> dict:
    p = _sources_json_path(nb_root)
    if not p.exists():
        return {"sources": []}
    return json.loads(p.read_text(encoding="utf-8"))


def _save_sources(nb_root: Path, data: dict) -> None:
    p = _sources_json_path(nb_root)
    p.write_text(json.dumps(data, indent=2), encoding="utf-8")


def _upsert_source(nb_root: Path, meta: SourceMeta) -> None:
    data = _load_sources(nb_root)
    items = data.get("sources", [])
    for i, it in enumerate(items):
        if it.get("id") == meta.id:
            items[i] = asdict(meta)
            break
    else:
        items.append(asdict(meta))
    data["sources"] = items
    _save_sources(nb_root, data)


def list_sources(username: str, notebook_id: str) -> List[SourceMeta]:
    nb_root = notebook_root(username, notebook_id)
    data = _load_sources(nb_root)
    out: List[SourceMeta] = []
    for it in data.get("sources", []):
        out.append(
            SourceMeta(
                id=str(it["id"]),
                type=str(it.get("type", "unknown")),
                enabled=bool(it.get("enabled", True)),
                raw_path=str(it.get("raw_path", "")),
                text_path=str(it.get("text_path", "")),
            )
        )
    return out


def set_source_enabled(username: str, notebook_id: str, source_id: str, enabled: bool) -> None:
    nb_root = notebook_root(username, notebook_id)
    data = _load_sources(nb_root)
    items = data.get("sources", [])
    for it in items:
        if it.get("id") == source_id:
            it["enabled"] = bool(enabled)
            break
    data["sources"] = items
    _save_sources(nb_root, data)


def _extract_pdf_text(pdf_path: Path) -> str:
    reader = PdfReader(str(pdf_path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append(f"\n\n--- Page {i + 1} ---\n{text}")
    return "\n".join(pages)


def _extract_txt_text(txt_path: Path) -> str:
    return txt_path.read_text(encoding="utf-8", errors="ignore")


def _extract_pptx_text(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    chunks = []
    for si, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            chunks.append(f"\n\n--- Slide {si} ---\n" + "\n".join(slide_text))
    return "\n".join(chunks)


def ingest_files(*, username: str, notebook_id: str, files: List[Path]) -> List[str]:
    nb_root = notebook_root(username, notebook_id)
    raw_dir = nb_root / "files_raw"
    ext_dir = nb_root / "files_extracted"
    raw_dir.mkdir(parents=True, exist_ok=True)
    ext_dir.mkdir(parents=True, exist_ok=True)

    ingested: List[str] = []

    for f in files:
        src = Path(f)
        if not src.exists():
            continue

        suffix = src.suffix.lower()
        if suffix not in (".pdf", ".txt", ".pptx"):
            continue

        dest = raw_dir / src.name
        shutil.copy2(src, dest)

        if suffix == ".pdf":
            text = _extract_pdf_text(dest)
            typ = "pdf"
        elif suffix == ".pptx":
            text = _extract_pptx_text(dest)
            typ = "pptx"
        else:
            text = _extract_txt_text(dest)
            typ = "txt"

        out_txt = ext_dir / f"{src.stem}.txt"
        out_txt.write_text(text, encoding="utf-8")

        _upsert_source(
            nb_root,
            SourceMeta(
                id=src.name,
                type=typ,
                enabled=True,
                raw_path=str(Path("files_raw") / src.name),
                text_path=str(Path("files_extracted") / f"{src.stem}.txt"),
            ),
        )

        ingested.append(src.name)

    return ingested